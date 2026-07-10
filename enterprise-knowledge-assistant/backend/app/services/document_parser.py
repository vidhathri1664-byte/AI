from pathlib import Path

import fitz
from docx import Document as DocxDocument
from pptx import Presentation


class DocumentParsingError(Exception):
    """Raised when a supported document cannot be parsed."""


def extract_text_from_pdf(file_path: Path) -> str:
    text_parts: list[str] = []

    try:
        with fitz.open(file_path) as pdf:
            for page_number, page in enumerate(pdf, start=1):
                page_text = page.get_text("text").strip()

                if page_text:
                    text_parts.append(
                        f"[Page {page_number}]\n{page_text}"
                    )
    except Exception as exc:
        raise DocumentParsingError(
            f"Unable to parse PDF: {file_path.name}"
        ) from exc

    return "\n\n".join(text_parts)


def extract_text_from_docx(file_path: Path) -> str:
    text_parts: list[str] = []

    try:
        document = DocxDocument(file_path)

        for paragraph in document.paragraphs:
            paragraph_text = paragraph.text.strip()

            if paragraph_text:
                text_parts.append(paragraph_text)

        for table_index, table in enumerate(document.tables, start=1):
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip()
                    for cell in row.cells
                    if cell.text.strip()
                )

                if row_text:
                    text_parts.append(
                        f"[Table {table_index}] {row_text}"
                    )
    except Exception as exc:
        raise DocumentParsingError(
            f"Unable to parse DOCX: {file_path.name}"
        ) from exc

    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_path: Path) -> str:
    text_parts: list[str] = []

    try:
        presentation = Presentation(file_path)

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            slide_parts: list[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = shape.text.strip()

                    if shape_text:
                        slide_parts.append(shape_text)

            if slide_parts:
                text_parts.append(
                    f"[Slide {slide_number}]\n"
                    + "\n".join(slide_parts)
                )
    except Exception as exc:
        raise DocumentParsingError(
            f"Unable to parse PPTX: {file_path.name}"
        ) from exc

    return "\n\n".join(text_parts)


def extract_text(file_path: str | Path) -> str:
    path = Path(file_path)

    if not path.exists():
        raise DocumentParsingError(
            f"Document does not exist: {path}"
        )

    extension = path.suffix.lower()

    parsers = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
    }

    parser = parsers.get(extension)

    if parser is None:
        raise DocumentParsingError(
            f"Unsupported document type: {extension}"
        )

    extracted_text = parser(path).strip()

    if not extracted_text:
        raise DocumentParsingError(
            f"No readable text found in {path.name}"
        )

    return extracted_text
